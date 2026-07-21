import os
import subprocess
import sys
import json
from pathlib import Path
import socket
from dotenv import load_dotenv
import convertapi

from PyQt6.QtCore import Qt, QMimeData, QUrl
from PyQt6.QtGui import QDrag
from PyQt6.QtWidgets import (
    QDialog,
    QVBoxLayout,
    QHBoxLayout,
    QLabel,
    QPushButton,
    QListWidget,
    QListWidgetItem,
    QFileDialog,
    QMessageBox,
    QComboBox
)

load_dotenv()
convertapi.api_secret = os.getenv("CONVERTAPI_SECRET")



class DraggableFileList(QListWidget):
    def __init__(self, file_basket, parent=None):
        super().__init__(parent)

        self.file_basket = file_basket
        self.setDragEnabled(True)
        self.setSelectionMode(QListWidget.SelectionMode.SingleSelection)
        self.setAcceptDrops(True)
        self.setDropIndicatorShown(True)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dragMoveEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dropEvent(self, event):
        urls = event.mimeData().urls()
        added_count = 0

        for url in urls:
            file_path = url.toLocalFile()

            if file_path:
                self.file_basket.add_file(file_path)
                added_count += 1

        if added_count > 0:
            parent_dialog = self.window()
            if hasattr(parent_dialog, "refresh_list"):
                parent_dialog.refresh_list()

        event.acceptProposedAction()

    def startDrag(self, supported_actions):
        selected_item = self.currentItem()

        if selected_item is None:
            return

        file_path = selected_item.data(Qt.ItemDataRole.UserRole)

        if not file_path or not os.path.exists(file_path):
            if file_path:
                self.file_basket.remove_file(file_path)
                self.takeItem(self.row(selected_item))
            return

        mime_data = QMimeData()
        mime_data.setUrls([QUrl.fromLocalFile(file_path)])

        drag = QDrag(self)
        drag.setMimeData(mime_data)

        result = drag.exec(Qt.DropAction.CopyAction | Qt.DropAction.MoveAction)

        if result != Qt.DropAction.IgnoreAction:
            self.file_basket.remove_file(file_path)
            self.takeItem(self.row(selected_item))


class FileBasketDialog(QDialog):
    def __init__(self, file_basket, parent=None):
        super().__init__(parent)

        self.file_basket = file_basket
        self.setWindowTitle("File Basket")
        self.setFixedSize(540, 480)

        # İŞTE SİHİR BURADA: Ask AI ve Task Manager'daki gibi Hayalet'in dinamik temasını çekiyoruz!
        if parent is not None and hasattr(parent, "get_dialog_stylesheet"):
            self.setStyleSheet(parent.get_dialog_stylesheet())

        if self.screen():
            screen_geom = self.screen().availableGeometry()
            x = int((screen_geom.width() - self.width()) / 2)
            y = int((screen_geom.height() - self.height()) / 2)
            self.move(x, y)

        layout = QVBoxLayout()
        layout.setContentsMargins(22, 20, 22, 20)
        layout.setSpacing(10)

        title_label = QLabel("FILE BASKET")
        title_label.setObjectName("dialogTitle")
        title_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

        info_label = QLabel(
            "Files held by ghost.\n"
            "Drag a file from this list back to Finder/Desktop.\n"
            "Select a PDF, click Convert To, then choose format and save location."
        )
        info_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

        section_label = QLabel("Dropped files")

        self.file_list = DraggableFileList(self.file_basket)
        self.refresh_list()

        top_button_layout = QHBoxLayout()
        top_button_layout.setSpacing(10)

        add_button = QPushButton("＋ ADD FILE")
        open_button = QPushButton("OPEN")
        remove_button = QPushButton("REMOVE")

        add_button.clicked.connect(self.add_file)
        open_button.clicked.connect(self.open_selected_file)
        remove_button.clicked.connect(self.remove_selected_file)

        top_button_layout.addWidget(add_button)
        top_button_layout.addWidget(open_button)
        top_button_layout.addWidget(remove_button)

        converter_button_layout = QHBoxLayout()
        converter_button_layout.setSpacing(10)

        convert_to_button = QPushButton("CONVERT TO")
        convert_to_button.clicked.connect(self.convert_selected_file)

        converter_button_layout.addWidget(convert_to_button)

        bottom_button_layout = QHBoxLayout()
        bottom_button_layout.setSpacing(10)

        clear_button = QPushButton("CLEAR BASKET")
        clear_button.setObjectName("dangerButton")
        close_button = QPushButton("CLOSE")
        close_button.setObjectName("dangerButton")

        clear_button.clicked.connect(self.clear_basket)
        close_button.clicked.connect(self.close)

        bottom_button_layout.addWidget(clear_button)
        bottom_button_layout.addWidget(close_button)

        layout.addWidget(title_label)
        layout.addWidget(info_label)
        layout.addSpacing(4)
        layout.addWidget(section_label)
        layout.addWidget(self.file_list)
        layout.addLayout(top_button_layout)
        layout.addLayout(converter_button_layout)
        layout.addLayout(bottom_button_layout)

        self.setLayout(layout)

    def show_info(self, title, message):
        info_box = QMessageBox(self)
        info_box.setWindowTitle(title)
        info_box.setText(message)
        info_box.setIcon(QMessageBox.Icon.Information)
        # Popup ekranı da ana temayı miras alıyor
        info_box.setStyleSheet(self.styleSheet())
        info_box.exec()

    def show_warning(self, title, message):
        warning_box = QMessageBox(self)
        warning_box.setWindowTitle(title)
        warning_box.setText(message)
        warning_box.setIcon(QMessageBox.Icon.Warning)
        # Popup ekranı da ana temayı miras alıyor
        warning_box.setStyleSheet(self.styleSheet())
        warning_box.exec()

    def add_file(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Choose a file",
            "",
            "All Files (*)"
        )

        if file_path:
            self.file_basket.add_file(file_path)
            self.refresh_list()

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
        else:
            event.ignore()

    def dropEvent(self, event):
        urls = event.mimeData().urls()
        added_count = 0

        for url in urls:
            file_path = url.toLocalFile()

            if file_path:
                self.file_basket.add_file(file_path)
                added_count += 1

        if added_count > 0:
            self.refresh_list()

        event.acceptProposedAction()

    def open_selected_file(self):
        selected_item = self.file_list.currentItem()

        if selected_item is None:
            self.show_warning("No Selection", "Please select a file.")
            return

        file_path = selected_item.data(Qt.ItemDataRole.UserRole)

        if file_path is None:
            self.show_warning("File Not Found", "This file could not be found.")
            return

        self.file_basket.open_file(file_path)

    def get_selected_file_path(self):
        selected_item = self.file_list.currentItem()

        if selected_item is None:
            self.show_warning("No Selection", "Please select a file first.")
            return None

        file_path = selected_item.data(Qt.ItemDataRole.UserRole)

        if not file_path or not os.path.exists(file_path):
            self.show_warning("File Not Found", "This file could not be found.")
            return None

        return file_path

    def ask_convert_format(self, convert_options):
        format_dialog = QDialog(self)
        format_dialog.setWindowTitle("Convert To")
        format_dialog.setFixedSize(340, 190)

        # Çeviri seçme ekranı da ana temayı miras alıyor
        format_dialog.setStyleSheet(self.styleSheet())

        layout = QVBoxLayout()
        layout.setContentsMargins(22, 18, 22, 18)
        layout.setSpacing(12)

        title_label = QLabel("CONVERT TO")
        title_label.setObjectName("dialogTitle")
        title_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

        hint_label = QLabel("Choose output format")

        format_box = QComboBox()
        format_box.addItems(convert_options)

        button_layout = QHBoxLayout()
        button_layout.setSpacing(10)

        convert_button = QPushButton("CONVERT")

        cancel_button = QPushButton("CANCEL")
        cancel_button.setObjectName("dangerButton")

        selected = {"value": None}

        def accept_format():
            selected["value"] = format_box.currentText()
            format_dialog.accept()

        convert_button.clicked.connect(accept_format)
        cancel_button.clicked.connect(format_dialog.reject)

        button_layout.addWidget(convert_button)
        button_layout.addWidget(cancel_button)

        layout.addWidget(title_label)
        layout.addWidget(hint_label)
        layout.addWidget(format_box)
        layout.addLayout(button_layout)

        format_dialog.setLayout(layout)

        if format_dialog.exec() == QDialog.DialogCode.Accepted:
            return selected["value"]

        return None

    def convert_selected_file(self):
        file_path = self.get_selected_file_path()

        if file_path is None:
            return

        source_path = Path(file_path)
        extension = source_path.suffix.lower()

        convert_options_by_type = {
            ".pdf": [
                "Word document (.docx)",
                "PowerPoint presentation (.pptx)",
                "Text file (.txt)"
            ],
            ".docx": [
                "PDF file (.pdf)",
                "Text file (.txt)"
            ],
            ".pptx": [
                "PDF file (.pdf)",
                "Text file (.txt)"
            ],
            ".txt": [
                "Word document (.docx)",
                "PDF file (.pdf)"
            ]
        }

        if extension not in convert_options_by_type:
            self.show_warning(
                "Unsupported File Type",
                "This file type cannot be converted yet. Try PDF, DOCX, PPTX, or TXT."
            )
            return

        convert_options = convert_options_by_type[extension]
        selected_format = self.ask_convert_format(convert_options)

        if selected_format is None:
            return

        if selected_format == "Word document (.docx)":
            output_extension = ".docx"
            save_filter = "Word Document (*.docx)"
            save_title = "Save converted Word file"

        elif selected_format == "PowerPoint presentation (.pptx)":
            output_extension = ".pptx"
            save_filter = "PowerPoint Presentation (*.pptx)"
            save_title = "Save converted PowerPoint file"

        elif selected_format == "PDF file (.pdf)":
            output_extension = ".pdf"
            save_filter = "PDF File (*.pdf)"
            save_title = "Save converted PDF file"

        else:
            output_extension = ".txt"
            save_filter = "Text File (*.txt)"
            save_title = "Save converted text file"

        default_output = str(source_path.with_suffix(output_extension))

        output_path, _ = QFileDialog.getSaveFileName(
            self,
            save_title,
            default_output,
            save_filter
        )

        if not output_path:
            return

        if not output_path.lower().endswith(output_extension):
            output_path += output_extension

        if extension == ".pdf" and output_extension == ".docx":
            ok, message = self.file_basket.pdf_to_docx(file_path, output_path)

        elif extension == ".pdf" and output_extension == ".pptx":
            ok, message = self.file_basket.pdf_to_pptx(file_path, output_path)

        elif extension == ".pdf" and output_extension == ".txt":
            ok, message = self.file_basket.pdf_to_txt(file_path, output_path)

        elif extension == ".docx" and output_extension == ".txt":
            ok, message = self.file_basket.docx_to_txt(file_path, output_path)

        elif extension == ".docx" and output_extension == ".pdf":
            ok, message = self.file_basket.docx_to_pdf(file_path, output_path)

        elif extension == ".pptx" and output_extension == ".txt":
            ok, message = self.file_basket.pptx_to_txt(file_path, output_path)

        elif extension == ".pptx" and output_extension == ".pdf":
            ok, message = self.file_basket.pptx_to_pdf(file_path, output_path)

        elif extension == ".txt" and output_extension == ".docx":
            ok, message = self.file_basket.txt_to_docx(file_path, output_path)

        elif extension == ".txt" and output_extension == ".pdf":
            ok, message = self.file_basket.txt_to_pdf(file_path, output_path)

        else:
            ok, message = False, "This conversion is not available yet."

        if ok:
            self.file_basket.add_file(output_path)
            self.refresh_list()
            self.show_info("Converted", message)
        else:
            self.show_warning("Conversion Error", message)

    def remove_selected_file(self):
        selected_item = self.file_list.currentItem()

        if selected_item is None:
            self.show_warning("No Selection", "Please select a file.")
            return

        file_path = selected_item.data(Qt.ItemDataRole.UserRole)
        self.file_basket.remove_file(file_path)
        self.refresh_list()

    def clear_basket(self):
        self.file_basket.clear_files()
        self.refresh_list()

    def refresh_list(self):
        self.file_list.clear()

        files = self.file_basket.get_files()

        if not files:
            empty_item = QListWidgetItem("Basket is empty.")
            empty_item.setFlags(empty_item.flags() & ~Qt.ItemFlag.ItemIsSelectable)
            self.file_list.addItem(empty_item)
            return

        for file_path in files:
            clean_path = file_path.rstrip("/")
            file_name = os.path.basename(clean_path)

            if os.path.isdir(file_path):
                display_name = f"📁  {file_name}"
            else:
                display_name = f"📄  {file_name}"

            item = QListWidgetItem(display_name)
            item.setData(Qt.ItemDataRole.UserRole, file_path)

            self.file_list.addItem(item)

class FileBasket:
    def __init__(self, ghost):
        self.ghost = ghost

        self.data_dir = "user_data"
        os.makedirs(self.data_dir, exist_ok=True)

        self.data_file = os.path.join(self.data_dir, "file_basket.json")
        self.files = self.load_data()

    def load_data(self):
        if os.path.exists(self.data_file):
            try:
                with open(self.data_file, "r", encoding="utf-8") as f:
                    saved_files = json.load(f)
                    valid_files = [path for path in saved_files if os.path.exists(path)]
                    return valid_files
            except Exception as error:
                print("Sepet verisi yüklenemedi:", error)
        return []

    def save_data(self):
        try:
            with open(self.data_file, "w", encoding="utf-8") as f:
                json.dump(self.files, f, indent=4)
        except Exception as error:
            print("Sepet verisi kaydedilemedi:", error)

    def open_dialog(self):
        dialog = FileBasketDialog(self, self.ghost)
        dialog.exec()

    def add_file(self, file_path):
        if not os.path.exists(file_path):
            self.ghost.speech_bubble.show_message("File does not exist.", 3000)
            return

        if file_path in self.files:
            self.ghost.speech_bubble.show_message("This file is already in basket.", 3000)
            return

        self.files.append(file_path)
        self.save_data()

        file_name = os.path.basename(file_path)
        self.ghost.speech_bubble.show_message(f"I got it: {file_name}", 3500)

    def get_files(self):
        return self.files

    def remove_file(self, file_path):
        if file_path in self.files:
            self.files.remove(file_path)
            self.save_data()

            file_name = os.path.basename(file_path)
            self.ghost.speech_bubble.show_message(f"Released: {file_name}", 3000)

    def clear_files(self):
        self.files.clear()
        self.save_data()
        self.ghost.speech_bubble.show_message("File basket cleared.", 3000)

    # --- YARDIMCI FONKSİYONLAR ---
    def check_internet(self):
        """İnternet bağlantısı olup olmadığını kontrol eder."""
        try:
            socket.create_connection(("1.1.1.1", 53), timeout=2)
            return True
        except OSError:
            return False

    def extract_pdf_text(self, pdf_path):
        try:
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            text_parts = []
            for page_number, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {page_number} ---\n{page_text.strip()}")
            full_text = "\n\n".join(text_parts).strip()
            if not full_text:
                return None, "No selectable text found. This may be a scanned PDF and needs OCR."
            return full_text, None
        except Exception as error:
            return None, f"Could not read PDF:\n{error}"

    # --- DÖNÜŞTÜRME FONKSİYONLARI (HİBRİT YAPI) ---

    def pdf_to_docx(self, pdf_path, output_path):
        if self.check_internet() and convertapi.api_secret:
            try:
                self.ghost.speech_bubble.show_message("Converting via Cloud...", 3000)
                result = convertapi.convert('docx', {'File': pdf_path}, from_format='pdf')
                result.save_files(output_path)
                return True, f"Saved via Cloud API as:\n{output_path}"
            except Exception as e:
                print("Cloud failed, trying local:", e)

        # B Planı: Yerel Dönüşüm
        try:
            self.ghost.speech_bubble.show_message("Converting locally...", 3000)
            from pdf2docx import Converter
            cv = Converter(pdf_path)
            cv.convert(output_path)
            cv.close()
            return True, f"Saved locally as:\n{output_path}"
        except Exception as error:
            return False, f"Could not convert PDF:\n{error}"

    def pdf_to_pptx(self, pdf_path, output_path):
        if self.check_internet() and convertapi.api_secret:
            try:
                self.ghost.speech_bubble.show_message("Converting via Cloud...", 3000)
                result = convertapi.convert('pptx', {'File': pdf_path}, from_format='pdf')
                result.save_files(output_path)
                return True, f"Saved via Cloud API as:\n{output_path}"
            except Exception as e:
                print("Cloud failed:", e)
        return False, "Cloud API quota reached or no internet. Local PDF to PPTX is not fully supported yet."

    def docx_to_pdf(self, docx_path, output_path):
        if self.check_internet() and convertapi.api_secret:
            try:
                self.ghost.speech_bubble.show_message("Converting via Cloud...", 3000)
                result = convertapi.convert('pdf', {'File': docx_path}, from_format='docx')
                result.save_files(output_path)
                return True, f"Saved via Cloud API as:\n{output_path}"
            except Exception as e:
                print("Cloud failed, trying local:", e)

        # B Planı: Yerel Dönüşüm
        try:
            self.ghost.speech_bubble.show_message("Converting locally...", 3000)
            from docx2pdf import convert
            convert(docx_path, output_path)
            return True, f"Saved locally as:\n{output_path}"
        except Exception as error:
            return False, f"Ensure MS Word is installed.\nError: {error}"

    def pptx_to_pdf(self, pptx_path, output_path):
        if self.check_internet() and convertapi.api_secret:
            try:
                self.ghost.speech_bubble.show_message("Converting via Cloud...", 3000)
                result = convertapi.convert('pdf', {'File': pptx_path}, from_format='pptx')
                result.save_files(output_path)
                return True, f"Saved via Cloud API as:\n{output_path}"
            except Exception as e:
                return False, f"Cloud conversion failed:\n{e}"
        return False, "No internet or API key missing for PPTX to PDF."

    # --- STANDART METİN DÖNÜŞÜMLERİ (YEREL) ---

    def pdf_to_txt(self, pdf_path, output_path):
        text, error = self.extract_pdf_text(pdf_path)
        if error: return False, error
        try:
            with open(output_path, "w", encoding="utf-8") as file:
                file.write(text)
            return True, f"Saved as:\n{output_path}"
        except Exception as error:
            return False, str(error)

    def docx_to_txt(self, docx_path, output_path):
        try:
            from docx import Document
            document = Document(docx_path)
            text = "\n".join([p.text.strip() for p in document.paragraphs if p.text.strip()])
            with open(output_path, "w", encoding="utf-8") as file:
                file.write(text)
            return True, f"Saved as:\n{output_path}"
        except Exception as error:
            return False, str(error)

    def pptx_to_txt(self, pptx_path, output_path):
        try:
            from pptx import Presentation
            presentation = Presentation(pptx_path)
            slide_texts = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    slide_texts.append(f"--- Slide {slide_index} ---\n" + "\n".join(texts))
            with open(output_path, "w", encoding="utf-8") as file:
                file.write("\n\n".join(slide_texts).strip())
            return True, f"Saved as:\n{output_path}"
        except Exception as error:
            return False, str(error)

    def txt_to_docx(self, txt_path, output_path):
        try:
            from docx import Document
            with open(txt_path, "r", encoding="utf-8") as file:
                text = file.read()
            document = Document()
            document.add_heading(Path(txt_path).stem, level=1)
            for block in text.split("\n\n"):
                if block.strip(): document.add_paragraph(block.strip())
            document.save(output_path)
            return True, f"Saved as:\n{output_path}"
        except Exception as error:
            return False, str(error)

    def txt_to_pdf(self, txt_path, output_path):
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            with open(txt_path, "r", encoding="utf-8") as file:
                text = file.read()
            pdf = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            x, y = 50, height - 50
            for line in text.splitlines():
                if not line.strip():
                    y -= 16
                    continue
                pdf.drawString(x, y, line[:90])
                y -= 16
                if y < 50:
                    pdf.showPage()
                    y = height - 50
            pdf.save()
            return True, f"Saved as:\n{output_path}"
        except Exception as error:
            return False, str(error)

    def open_file(self, file_path):
        if not os.path.exists(file_path):
            self.ghost.speech_bubble.show_message("File not found.", 3000)
            return
        try:
            if sys.platform == "darwin":
                subprocess.call(["open", file_path])
            elif sys.platform == "win32":
                os.startfile(file_path)
            else:
                subprocess.call(["xdg-open", file_path])
        except Exception as error:
            print("Open Error:", error)