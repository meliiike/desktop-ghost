import os

from PyQt6.QtWidgets import (
    QDialog,
    QWidget,
    QLabel,
    QVBoxLayout,
    QHBoxLayout,
    QTextEdit,
    QLineEdit,
    QPushButton,
    QFileDialog,
    QMessageBox,
    QComboBox,
)

from PyQt6.QtCore import Qt, QEvent
from .gemini_ai import GeminiAI


class AttachFileDialog(QDialog):
    def __init__(self, ask_dialog, parent=None):
        super().__init__(parent)

        self.ask_dialog = ask_dialog
        self.ghost = ask_dialog.ghost

        self.setWindowTitle("Attach File")
        self.setFixedSize(340, 168)
        self.setWindowFlags(
            Qt.WindowType.FramelessWindowHint |
            Qt.WindowType.Popup
        )
        self.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground)

        self.setup_ui()
        self.refresh_basket_files()

    def setup_ui(self):
        self.setStyleSheet("""
            QDialog {
                background-color: transparent;
            }

            QWidget#attachCard {
                background-color: rgba(255, 253, 235, 250);
                border: 3px solid #8B7A5A;
                border-radius: 18px;
            }

            QLabel#attachTitle {
                color: #4A3B24;
                font-size: 13px;
                font-weight: bold;
                padding: 2px;
            }

            QLabel#attachHint {
                color: #5B4A31;
                font-size: 10px;
                font-weight: bold;
            }

            QComboBox {
                background-color: #FFFFF5;
                color: #4A3B24;
                border: 2px solid rgba(139, 122, 90, 150);
                border-radius: 10px;
                padding: 7px;
                font-size: 11px;
                font-weight: bold;
            }

            QComboBox:hover {
                border: 2px solid #8B7A5A;
            }

            QComboBox QAbstractItemView {
                background-color: #FFFDEB;
                color: #4A3B24;
                selection-background-color: #F1D9AA;
                selection-color: #3A2A15;
                border: 2px solid #8B7A5A;
                outline: none;
            }

            QComboBox::drop-down {
                border: none;
                width: 22px;
            }

            QComboBox::down-arrow {
                image: none;
            }

            QPushButton {
                background-color: #FFF1CD;
                color: #4A3B24;
                border: 2px solid #8B7A5A;
                border-radius: 10px;
                padding: 7px 9px;
                font-size: 11px;
                font-weight: bold;
            }

            QPushButton:hover {
                background-color: #F5E8C8;
                border: 2px solid #A88955;
            }

            QPushButton:pressed {
                background-color: #EBD2AA;
                padding-top: 8px;
            }

            QPushButton#cancelButton {
                color: #8A3D2B;
                border: 2px solid #A45A3A;
            }
        """)

        outer_layout = QVBoxLayout(self)
        outer_layout.setContentsMargins(0, 0, 0, 0)

        card = QWidget()
        card.setObjectName("attachCard")

        layout = QVBoxLayout(card)
        layout.setContentsMargins(15, 13, 15, 13)
        layout.setSpacing(8)

        title = QLabel("📎 ATTACH FILE")
        title.setObjectName("attachTitle")
        title.setAlignment(Qt.AlignmentFlag.AlignCenter)

        hint = QLabel("Pick a basket file or browse from device.")
        hint.setObjectName("attachHint")
        hint.setAlignment(Qt.AlignmentFlag.AlignCenter)

        self.basket_file_box = QComboBox()

        button_row = QHBoxLayout()
        button_row.setSpacing(7)

        use_button = QPushButton("USE SELECTED")
        browse_button = QPushButton("BROWSE")
        cancel_button = QPushButton("CANCEL")
        cancel_button.setObjectName("cancelButton")

        use_button.clicked.connect(self.use_selected_basket_file)
        browse_button.clicked.connect(self.choose_from_computer)
        cancel_button.clicked.connect(self.close)

        for button in (use_button, browse_button, cancel_button):
            button.setAutoDefault(False)
            button.setDefault(False)

        button_row.addWidget(use_button)
        button_row.addWidget(browse_button)
        button_row.addWidget(cancel_button)

        layout.addWidget(title)
        layout.addWidget(hint)
        layout.addWidget(self.basket_file_box)
        layout.addLayout(button_row)

        outer_layout.addWidget(card)

    def refresh_basket_files(self):
        self.basket_file_box.clear()
        self.basket_file_box.addItem("No basket file selected", None)

        if not hasattr(self.ghost, "file_basket"):
            self.basket_file_box.addItem("File Basket is not available", None)
            return

        files = self.ghost.file_basket.get_files()

        for file_path in files:
            clean_path = file_path.rstrip("/")
            file_name = os.path.basename(clean_path)

            if os.path.isdir(file_path):
                display_name = f"📁 {file_name}"
            else:
                display_name = f"📄 {file_name}"

            self.basket_file_box.addItem(display_name, file_path)

    def use_selected_basket_file(self):
        file_path = self.basket_file_box.currentData()

        if not file_path:
            self.ask_dialog.show_warning("No File", "Please choose a file from File Basket.")
            return

        if not os.path.exists(file_path):
            self.ask_dialog.show_warning("File Not Found", "This file does not exist anymore.")
            self.refresh_basket_files()
            return

        self.ask_dialog.set_selected_file(file_path, "Selected from basket")
        self.close()

    def choose_from_computer(self):
        # Use the main Ask AI dialog as parent.
        # This avoids the small popup blocking/closing the Windows file picker.
        file_path, _ = QFileDialog.getOpenFileName(
            self.ask_dialog,
            "Choose a file",
            "",
            "Documents (*.docx *.pdf *.pptx *.txt);;All Files (*)",
            options=QFileDialog.Option.DontUseNativeDialog
        )

        if file_path:
            self.ask_dialog.set_selected_file(file_path, "Selected file")
            self.close()


class AskAIDialog(QDialog):
    def __init__(self, ghost, parent=None):
        super().__init__(parent)

        self.ghost = ghost
        self.selected_file = None
        self.ai = None
        self.drag_position = None

        self.setWindowTitle("Ask AI")
        self.setFixedSize(420, 460)

        self.setWindowFlags(
            Qt.WindowType.FramelessWindowHint |
            Qt.WindowType.Dialog
        )

        self.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground)

        self.setup_ui()
        self.append_ai_message("Hi! Ask me something.")

    def setup_ui(self):
        self.setStyleSheet("""
            QDialog {
                background-color: transparent;
            }

            QWidget#aiCard {
                background-color: rgba(255, 253, 235, 245);
                border: 3px solid #8B7A5A;
                border-radius: 22px;
            }

            QLabel#title {
                background-color: rgba(255, 241, 205, 230);
                color: #4A3B24;
                border: 2px solid #8B7A5A;
                border-radius: 14px;
                padding: 8px 18px;
                font-size: 18px;
                font-weight: bold;
                letter-spacing: 2px;
            }

            QLabel#subtitle, QLabel#selectedFileLabel {
                color: #4A3B24;
                font-size: 11px;
                font-weight: bold;
            }

            QTextEdit {
                background-color: rgba(255, 255, 245, 230);
                color: #4A3B24;
                border: 2px solid rgba(139, 122, 90, 120);
                border-radius: 13px;
                padding: 9px;
                font-size: 12px;
            }

            QLineEdit {
                background-color: rgba(255, 255, 245, 235);
                color: #4A3B24;
                border: 2px solid #8B7A5A;
                border-radius: 11px;
                padding: 9px;
                font-size: 12px;
            }

            QLineEdit:focus {
                background-color: rgba(255, 255, 255, 245);
                border: 2px solid #A88955;
            }

            QPushButton {
                background-color: rgba(255, 241, 205, 230);
                color: #4A3B24;
                border: 2px solid #8B7A5A;
                border-radius: 11px;
                padding: 8px 10px;
                font-size: 11px;
                font-weight: bold;
            }

            QPushButton:hover {
                background-color: rgba(245, 232, 200, 230);
                border: 2px solid #A88955;
            }

            QPushButton:pressed {
                background-color: rgba(235, 210, 170, 230);
                padding-top: 9px;
            }

            QPushButton#attachButton {
                font-size: 15px;
                padding: 4px;
            }

            QPushButton#closeButton {
                color: #8A3D2B;
                border: 2px solid #A45A3A;
            }

            QPushButton#closeButton:hover {
                background-color: rgba(255, 220, 190, 230);
            }
        """)

        outer_layout = QVBoxLayout(self)
        outer_layout.setContentsMargins(0, 0, 0, 0)

        card = QWidget()
        card.setObjectName("aiCard")

        layout = QVBoxLayout(card)
        layout.setContentsMargins(20, 18, 20, 18)
        layout.setSpacing(8)

        title = QLabel("ASK AI")
        title.setObjectName("title")
        title.setAlignment(Qt.AlignmentFlag.AlignCenter)
        title.installEventFilter(self)

        subtitle = QLabel("Ask your AI assistant something.")
        subtitle.setObjectName("subtitle")
        subtitle.setAlignment(Qt.AlignmentFlag.AlignCenter)
        subtitle.installEventFilter(self)

        self.chat_area = QTextEdit()
        self.chat_area.setReadOnly(True)
        self.chat_area.setPlaceholderText("AI answers will appear here...")
        self.chat_area.setMinimumHeight(255)

        self.selected_file_label = QLabel("No file attached")
        self.selected_file_label.setObjectName("selectedFileLabel")
        self.selected_file_label.setWordWrap(True)

        self.input_box = QLineEdit()
        self.input_box.setPlaceholderText("Type your question...")
        self.input_box.returnPressed.connect(self.ask_ai)

        input_layout = QHBoxLayout()
        input_layout.setSpacing(7)

        attach_button = QPushButton("📎")
        attach_button.setObjectName("attachButton")
        attach_button.setFixedWidth(42)
        attach_button.clicked.connect(self.open_attach_dialog)
        attach_button.setAutoDefault(False)
        attach_button.setDefault(False)

        ask_button = QPushButton("ASK")
        ask_button.setFixedWidth(68)
        ask_button.clicked.connect(self.ask_ai)
        ask_button.setAutoDefault(False)
        ask_button.setDefault(False)

        input_layout.addWidget(attach_button)
        input_layout.addWidget(self.input_box)
        input_layout.addWidget(ask_button)

        close_button = QPushButton("✕ CLOSE")
        close_button.setObjectName("closeButton")
        close_button.clicked.connect(self.close)
        close_button.setAutoDefault(False)
        close_button.setDefault(False)

        layout.addWidget(title)
        layout.addWidget(subtitle)
        layout.addWidget(self.chat_area)
        layout.addWidget(self.selected_file_label)
        layout.addLayout(input_layout)
        layout.addWidget(close_button)

        outer_layout.addWidget(card)


    def eventFilter(self, watched, event):
        if event.type() == QEvent.Type.MouseButtonPress:
            if event.button() == Qt.MouseButton.LeftButton:
                self.drag_position = (
                    event.globalPosition().toPoint()
                    - self.frameGeometry().topLeft()
                )
                return True

        if event.type() == QEvent.Type.MouseMove:
            if self.drag_position is not None and event.buttons() & Qt.MouseButton.LeftButton:
                self.move(event.globalPosition().toPoint() - self.drag_position)
                return True

        if event.type() == QEvent.Type.MouseButtonRelease:
            self.drag_position = None
            return True

        return super().eventFilter(watched, event)

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.drag_position = (
                event.globalPosition().toPoint()
                - self.frameGeometry().topLeft()
            )
            event.accept()
        else:
            super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        if self.drag_position is not None and event.buttons() & Qt.MouseButton.LeftButton:
            self.move(event.globalPosition().toPoint() - self.drag_position)
            event.accept()
        else:
            super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        self.drag_position = None
        super().mouseReleaseEvent(event)


    def append_user_message(self, message):
        safe_message = self.escape_html(message)
        self.chat_area.append(
            f'<p style="margin:6px 0;"><b style="color:#1D8CB3;">You:</b> '
            f'<span style="color:#2F3A3D;">{safe_message}</span></p>'
        )

    def append_ai_message(self, message):
        safe_message = self.escape_html(message).replace("\n", "<br>")
        self.chat_area.append(
            f'<p style="margin:6px 0;"><b style="color:#3D9A4B;">AI:</b> '
            f'<span style="color:#2F3A3D;">{safe_message}</span></p>'
        )

    def append_system_message(self, message):
        safe_message = self.escape_html(message)
        self.chat_area.append(
            f'<p style="margin:6px 0;"><b style="color:#8A5A00;">File:</b> '
            f'<span style="color:#5B4A31;">{safe_message}</span></p>'
        )

    def escape_html(self, text):
        return (
            str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
        )

    def format_ai_error(self, error):
        error_text = str(error)

        if "getaddrinfo failed" in error_text or "11002" in error_text:
            return (
                "AI connection failed.\n\n"
                "This usually means internet/DNS connection could not reach Gemini.\n"
                "Check Wi-Fi/VPN/firewall and run again."
            )

        return f"Could not get AI response:\n{error_text}"

    def open_attach_dialog(self):
        dialog = AttachFileDialog(self, self)

        # Açılan küçük kutu Ask AI penceresinin alt kısmında dursun.
        dialog_x = self.x() + 18
        dialog_y = self.y() + self.height() - dialog.height() - 65
        dialog.move(dialog_x, dialog_y)
        dialog.exec()

    def set_selected_file(self, file_path, source_label):
        self.selected_file = file_path
        file_name = os.path.basename(file_path)
        self.selected_file_label.setText(f"📎 {file_name}")
        self.append_system_message(f"{source_label}: {file_name}")
        self.ghost.speech_bubble.show_message(f"Attached: {file_name}", 2500)

    def read_text_file(self, file_path):
        encodings = ["utf-8", "utf-8-sig", "cp1254", "latin-1"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue

        raise ValueError("This text file could not be decoded.")

    def read_pdf_file(self, file_path):
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("PDF reading needs pypdf. Run: pip install pypdf")

        reader = PdfReader(file_path)
        pages = []

        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Page {index} ---\n{text}")

        return "\n\n".join(pages).strip()

    def read_docx_file(self, file_path):
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Word reading needs python-docx. Run: pip install python-docx")

        document = Document(file_path)
        paragraphs = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)

        return "\n".join(paragraphs).strip()

    def read_pptx_file(self, file_path):
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("PowerPoint reading needs python-pptx. Run: pip install python-pptx")

        presentation = Presentation(file_path)
        slide_texts = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)

            if texts:
                slide_texts.append(f"--- Slide {slide_index} ---\n" + "\n".join(texts))

        return "\n\n".join(slide_texts).strip()

    def read_selected_file_content(self):
        if not self.selected_file:
            return None

        if not os.path.exists(self.selected_file):
            raise FileNotFoundError("Attached file was not found.")

        extension = os.path.splitext(self.selected_file)[1].lower()

        if extension in [".txt", ".py", ".md", ".csv", ".json", ".html", ".css", ".js"]:
            content = self.read_text_file(self.selected_file)
        elif extension == ".pdf":
            content = self.read_pdf_file(self.selected_file)
        elif extension == ".docx":
            content = self.read_docx_file(self.selected_file)
        elif extension == ".pptx":
            content = self.read_pptx_file(self.selected_file)
        else:
            raise ValueError("This file type is not readable yet. Try TXT, PDF, DOCX, PPTX, PY, MD, CSV, JSON.")

        content = content.strip()

        if not content:
            raise ValueError("I could not find readable text in this file.")

        max_chars = 12000
        if len(content) > max_chars:
            content = content[:max_chars] + "\n\n[File text was shortened because it is long.]"

        return content

    def ask_ai(self):
        user_text = self.input_box.text().strip()

        if not user_text:
            return

        self.append_user_message(user_text)
        self.input_box.clear()

        try:
            if self.ai is None:
                self.ai = GeminiAI()

            if self.selected_file:
                file_name = os.path.basename(self.selected_file)
                file_content = self.read_selected_file_content()

                user_text = (
                    f"The user attached this file: {file_name}\n\n"
                    f"File content:\n{file_content}\n\n"
                    f"User question: {user_text}"
                )

            answer = self.ai.ask(user_text)

            self.append_ai_message(answer)
            self.ghost.speech_bubble.show_message(answer[:38], 2500)


        except Exception as error:

            error_text = str(error)

            if "503" in error_text or "UNAVAILABLE" in error_text or "high demand" in error_text:
                self.add_ai_message(
                    "Gemini is busy right now. Please try again in a few seconds."
                )

                self.ghost.speech_bubble.show_message(
                    "AI is busy. Try again soon.",

                    3000
                )
                return

            if "getaddrinfo failed" in error_text:
                self.add_ai_message(
                    "I cannot connect to the internet right now. Please check your connection."
                )

                self.ghost.speech_bubble.show_message(
                    "Connection problem.",
                    3000
                )
                return

            self.add_ai_message(
                "Something went wrong while getting the AI response."
            )

            print("AI Error:", error)

    def show_warning(self, title, message):
        warning_box = QMessageBox(self)
        warning_box.setWindowTitle(title)
        warning_box.setText(message)
        warning_box.setIcon(QMessageBox.Icon.Warning)

        warning_box.setStyleSheet("""
            QMessageBox {
                background-color: #FFFDEB;
            }

            QLabel {
                color: #4A3B24;
                font-size: 12px;
                font-weight: bold;
            }

            QPushButton {
                background-color: #FFFFFF;
                color: #4A3B24;
                border: 1px solid #8B7A5A;
                border-radius: 9px;
                padding: 7px 14px;
                font-size: 12px;
                font-weight: bold;
                min-width: 70px;
            }

            QPushButton:hover {
                background-color: #FFF1CD;
            }
        """)

        warning_box.exec()
